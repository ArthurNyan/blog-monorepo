#!/usr/bin/env python3

from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "presentation"
GENERATED_DIR = OUTPUT_DIR / "assets-generated"
OUTPUT_PPTX = OUTPUT_DIR / "diploma-presentation.pptx"

TITLE = "Разработка системы управления контентом\nв сфере маркетинга на базе Strapi и Astro"
AUTHOR = "Нахатакян Артур Романович"
SUPERVISOR = (
    "кандидат пед. наук, доцент кафедры информационных технологий\n"
    "и электронного обучения Государев Илья Борисович"
)

BG = RGBColor(247, 245, 240)
TEXT = RGBColor(34, 39, 46)
MUTED = RGBColor(89, 99, 112)
ACCENT = RGBColor(21, 54, 93)
ACCENT_2 = RGBColor(193, 126, 50)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(219, 222, 227)
SOFT = RGBColor(234, 229, 221)


def ensure_dirs() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)


def convert_pdf_to_png(pdf_name: str) -> Path:
    pdf_path = ROOT / "assets" / "diagrams" / "plantuml" / "pdf" / pdf_name
    png_path = GENERATED_DIR / pdf_name.replace(".pdf", ".png")
    if png_path.exists() and png_path.stat().st_mtime >= pdf_path.stat().st_mtime:
        return png_path

    subprocess.run(
        [
            "pdftoppm",
            "-r",
            "220",
            "-png",
            "-singlefile",
            str(pdf_path),
            str(png_path.with_suffix("")),
        ],
        check=True,
    )
    return png_path


def build_assets() -> dict[str, Path]:
    return {
        "logo": ROOT / "assets" / "title-2025-logo.png",
        "home": ROOT / "assets" / "front" / "home-en-with-url.png",
        "vacancy": ROOT / "assets" / "front" / "frontend-vacancy.png",
        "page_editor": ROOT / "assets" / "strapi-images" / "page.png",
        "dynamic_zone": ROOT / "assets" / "strapi-images" / "home-page-dynamic-zone.png",
        "preview": ROOT / "assets" / "strapi-images" / "page-preview-desktop.png",
        "architecture": convert_pdf_to_png("cms-first-architecture.pdf"),
        "data_model": convert_pdf_to_png("cms-data-model.pdf"),
        "content_lifecycle": convert_pdf_to_png("cms-content-lifecycle.pdf"),
    }


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_band(slide, title: str, subtitle: str | None = None, slide_no: int | None = None) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.95))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.42), Inches(0.2), Inches(0.16), Inches(0.56))
    marker.fill.solid()
    marker.fill.fore_color.rgb = ACCENT_2
    marker.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.76), Inches(0.15), Inches(10.2), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.paragraphs[0].text = title
    title_frame.paragraphs[0].font.name = "Aptos Display"
    title_frame.paragraphs[0].font.size = Pt(26)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.78), Inches(0.52), Inches(10.3), Inches(0.2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.paragraphs[0].text = subtitle
        subtitle_frame.paragraphs[0].font.name = "Aptos"
        subtitle_frame.paragraphs[0].font.size = Pt(10.5)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(226, 231, 237)

    if slide_no is not None:
        slide_no_box = slide.shapes.add_textbox(Inches(12.3), Inches(0.18), Inches(0.55), Inches(0.35))
        paragraph = slide_no_box.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph.text = f"{slide_no}"
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(15)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(236, 241, 246)


def add_logo(slide, logo_path: Path, left: float, top: float, width: float) -> None:
    slide.shapes.add_picture(str(logo_path), Inches(left), Inches(top), width=Inches(width))


def add_text_block(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    paragraphs: list[dict],
    fill_rgb: RGBColor | None = None,
    margin: float = 0.14,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_rgb is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_rgb
        box.line.color.rgb = LINE
    else:
        box.fill.background()
        box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    tf.clear()
    for idx, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = spec["text"]
        p.level = spec.get("level", 0)
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(spec.get("space_after", 6))
        p.space_before = Pt(spec.get("space_before", 0))
        p.line_spacing = spec.get("line_spacing", 1.05)
        p.bullet = spec.get("bullet", False)
        font = p.font
        font.name = spec.get("font_name", "Aptos")
        font.size = Pt(spec.get("size", 18))
        font.bold = spec.get("bold", False)
        font.italic = spec.get("italic", False)
        font.color.rgb = spec.get("color", TEXT)


def add_card_title(slide, left: float, top: float, width: float, title: str) -> None:
    add_text_block(
        slide,
        left,
        top,
        width,
        0.28,
        [
            {
                "text": title,
                "size": 12,
                "bold": True,
                "color": ACCENT,
                "space_after": 0,
            }
        ],
    )


def add_bullet_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    bullets: list[str],
    title_color: RGBColor = ACCENT,
) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = LINE

    add_text_block(
        slide,
        left + 0.05,
        top + 0.04,
        width - 0.1,
        height - 0.08,
        [{"text": title, "size": 18, "bold": True, "color": title_color, "space_after": 8}]
        + [{"text": bullet, "size": 13.5, "bullet": True, "color": TEXT, "space_after": 3} for bullet in bullets],
    )


def add_image(slide, path: Path, left: float, top: float, width: float, height: float | None = None) -> None:
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=None if height is None else Inches(height))


def add_caption(slide, left: float, top: float, width: float, text: str) -> None:
    add_text_block(
        slide,
        left,
        top,
        width,
        0.28,
        [{"text": text, "size": 9.5, "color": MUTED, "italic": True, "space_after": 0}],
    )


def add_label_chip(slide, left: float, top: float, text: str, fill: RGBColor = SOFT, font: RGBColor = ACCENT) -> None:
    width = 0.44 + 0.075 * len(text)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    add_text_block(
        slide,
        left + 0.02,
        top + 0.01,
        width - 0.04,
        0.26,
        [{"text": text, "size": 10.5, "bold": True, "color": font, "align": PP_ALIGN.CENTER, "space_after": 0}],
        margin=0.02,
    )


def add_table_slide(slide, left: float, top: float, width: float, height: float, rows: list[tuple[str, str]]) -> None:
    table_shape = slide.shapes.add_table(len(rows) + 1, 2, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    table.columns[0].width = Inches(width * 0.42)
    table.columns[1].width = Inches(width * 0.58)

    headers = ["Что проверено", "Результат"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Aptos"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    for row_idx, (left_text, right_text) in enumerate(rows, start=1):
        for col_idx, cell_text in enumerate((left_text, right_text)):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if row_idx % 2 else RGBColor(250, 250, 252)
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Aptos"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT


def remove_empty_placeholders(slide) -> None:
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and not shape.has_text_frame:
            element = shape._element
            element.getparent().remove(element)


def add_title_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.45))
    hero.fill.solid()
    hero.fill.fore_color.rgb = ACCENT
    hero.line.fill.background()

    accent_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.9), Inches(13.333), Inches(0.6))
    accent_band.fill.solid()
    accent_band.fill.fore_color.rgb = ACCENT
    accent_band.line.fill.background()

    add_logo(slide, assets["logo"], 11.6, 0.18, 1.2)

    add_text_block(
        slide,
        0.72,
        0.34,
        10.2,
        0.7,
        [
            {
                "text": "ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА",
                "size": 12,
                "bold": True,
                "color": RGBColor(221, 228, 236),
                "space_after": 4,
            },
            {
                "text": "09.03.01 Информатика и вычислительная техника\nТехнологии разработки программного обеспечения",
                "size": 10.5,
                "color": RGBColor(230, 236, 242),
                "space_after": 0,
            },
        ],
    )

    add_text_block(
        slide,
        0.8,
        1.85,
        8.8,
        1.75,
        [
            {"text": TITLE, "size": 24, "bold": True, "font_name": "Aptos Display", "space_after": 10},
            {
                "text": "Тема сфокусирована на CMS-first архитектуре для маркетингового сайта: Strapi управляет контентом, Astro формирует публичную витрину.",
                "size": 14,
                "color": MUTED,
            },
        ],
    )

    add_bullet_card(
        slide,
        0.82,
        4.15,
        5.8,
        1.6,
        "Исполнитель",
        [
            AUTHOR,
            "4 курс, очная форма обучения",
        ],
        title_color=ACCENT,
    )
    add_bullet_card(
        slide,
        6.88,
        4.15,
        5.58,
        1.6,
        "Руководитель",
        [
            SUPERVISOR,
        ],
        title_color=ACCENT,
    )

    add_text_block(
        slide,
        0.84,
        6.98,
        12.0,
        0.24,
        [{"text": "Российский государственный педагогический университет им. А. И. Герцена · Санкт-Петербург · 2026", "size": 12, "bold": True, "color": RGBColor(255, 255, 255), "align": PP_ALIGN.CENTER, "space_after": 0}],
    )


def add_relevance_slide(prs: Presentation, assets: dict[str, Path], slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_band(slide, "Актуальность", "Почему для маркетингового сайта нужен CMS-first подход", slide_no)

    add_bullet_card(
        slide,
        0.72,
        1.28,
        6.0,
        4.8,
        "Проблема",
        [
            "Корпоративный маркетинговый сайт требует частых изменений страниц, кейсов, статей и вакансий.",
            "При кодо-центричном подходе тексты, SEO и локализация оказываются распределены между CMS, шаблонами и ручными правками.",
            "Редакторский процесс замедляется, потому что типовые изменения зависят от разработчика.",
            "Публикация должна быть предсказуемой: черновик, предпросмотр, публикация и обновление витрины.",
        ],
    )

    add_image(slide, assets["home"], 7.05, 1.4, 5.55)
    add_caption(slide, 7.08, 5.95, 5.35, "Публичная витрина должна обновляться из CMS, а не через ручную правку шаблонов.")

    add_label_chip(slide, 7.08, 6.42, "ru/en локализация")
    add_label_chip(slide, 8.82, 6.42, "preview mode")
    add_label_chip(slide, 10.48, 6.42, "SEO / sitemap")


def add_goal_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_band(slide, "Объект, предмет, цель и задачи", None, slide_no)

    add_bullet_card(
        slide,
        0.72,
        1.3,
        5.9,
        5.75,
        "Исследовательская рамка",
        [
            "Объект: процессы управления контентом корпоративного маркетингового сайта.",
            "Предмет: архитектурные и программные решения для CMS-first платформы на базе Strapi и Astro.",
            "Цель: разработать систему, позволяющую управлять страницами, статьями, кейсами и вакансиями без участия разработчика.",
        ],
    )

    add_bullet_card(
        slide,
        6.86,
        1.3,
        5.78,
        5.75,
        "Ключевые задачи",
        [
            "Обосновать выбор headless CMS и связки Strapi + Astro.",
            "Спроектировать требования, архитектуру и модель данных.",
            "Реализовать CMS-контур для pages, articles, projects, vacancies.",
            "Добавить ru/en, preview mode, SEO/Open Graph и sitemap.",
            "Организовать roles/permissions, webhook -> rebuild, Vercel и Docker.",
            "Проверить результат по сценариям публикации, локализации и безопасности.",
        ],
    )


def add_stack_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_band(slide, "Выбор подхода и технологий", None, slide_no)

    add_bullet_card(
        slide,
        0.72,
        1.45,
        3.95,
        4.9,
        "Headless CMS",
        [
            "Разделяет хранение контента и рендеринг витрины.",
            "Позволяет централизовать editor workflow и API-контракт.",
            "Подходит для мультиязычного маркетингового сайта с различными публичными разделами.",
        ],
        title_color=ACCENT_2,
    )
    add_bullet_card(
        slide,
        4.93,
        1.45,
        3.95,
        4.9,
        "Strapi 5",
        [
            "Гибкие content types, Dynamic Zone и media library.",
            "Встроенные i18n, RBAC и webhooks для редакторского контура.",
            "Подходит как CMS-ядро для pages, articles, projects и vacancies.",
        ],
    )
    add_bullet_card(
        slide,
        9.14,
        1.45,
        3.5,
        4.9,
        "Astro 6",
        [
            "Prerender-подход для предсказуемой отдачи контента.",
            "Server routes для preview и форменных сценариев.",
            "Удобная интеграция с sitemap и deployment на Vercel.",
        ],
    )

    add_text_block(
        slide,
        0.9,
        6.55,
        11.8,
        0.36,
        [
            {
                "text": "Итоговый выбор обоснован не абстрактно, а через задачу: передать управление маркетинговым контентом в CMS и оставить витрину легкой, статически собираемой и контролируемой.",
                "size": 13.5,
                "color": MUTED,
                "align": PP_ALIGN.CENTER,
                "space_after": 0,
            }
        ],
    )


def add_architecture_slide(prs: Presentation, assets: dict[str, Path], slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_band(slide, "Архитектура системы", "Monorepo, разделение CMS и витрины, публикационный контур", slide_no)

    add_image(slide, assets["architecture"], 0.8, 1.35, 8.1)
    add_caption(slide, 0.88, 6.1, 7.7, "Базовая схема: Strapi управляет контентом и ролями, Astro собирает публичные страницы, внешняя витрина обновляется через rebuild.")

    add_bullet_card(
        slide,
        9.15,
        1.52,
        3.55,
        4.55,
        "Ключевые компоненты",
        [
            "apps/cms: Strapi 5, content model, roles, webhooks.",
            "apps/front: Astro 6, публичные маршруты, preview и API routes.",
            "OpenAPI/typed contract: единая схема данных между CMS и frontend.",
            "Vercel + Docker: эксплуатационный контур системы.",
        ],
    )


def add_data_model_slide(prs: Presentation, assets: dict[str, Path], slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_band(slide, "Модель данных CMS", "Контентные сущности и общий SEO/локализационный слой", slide_no)

    add_image(slide, assets["data_model"], 0.78, 1.28, 8.65)
    add_caption(slide, 0.88, 5.43, 8.0, "Сущности ориентированы на предметную область: страницы, публикации, кейсы, вакансии и связанные справочники.")

    add_bullet_card(
        slide,
        9.62,
        1.4,
        3.0,
        4.75,
        "Что включает модель",
        [
            "core: global, home-page, page;",
            "content: article, project, vacancy;",
            "relations: author, industry, job-role;",
            "shared.seo для detail-страниц;",
            "локализуемые поля для storefront-core, articles и projects;",
            "форменные сущности lead-submission и vacancy-application.",
        ],
    )


def add_editor_slide(prs: Presentation, assets: dict[str, Path], slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_band(slide, "Редакторский контур", "Pages, Dynamic Zone и preview без публикации черновика", slide_no)

    add_text_block(slide, 0.78, 1.15, 5.7, 0.26, [{"text": "Редактирование страницы в Strapi", "size": 13, "bold": True, "color": ACCENT, "space_after": 0}])
    add_text_block(slide, 6.78, 1.15, 5.7, 0.26, [{"text": "Предпросмотр итоговой витрины", "size": 13, "bold": True, "color": ACCENT, "space_after": 0}])
    add_image(slide, assets["page_editor"], 0.78, 1.45, 5.62)
    add_image(slide, assets["preview"], 6.78, 1.45, 5.62)

    add_bullet_card(
        slide,
        0.82,
        5.75,
        12.0,
        1.0,
        "Практический результат",
        [
            "Сущность page стала самостоятельной контентной единицей со slug, SEO и блоками Dynamic Zone; preview mode открывает draft-версию через защищенный server-side сценарий без раскрытия черновиков в публичном API.",
        ],
        title_color=ACCENT_2,
    )


def add_frontend_slide(prs: Presentation, assets: dict[str, Path], slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_band(slide, "Публичная витрина и SEO-контур", "Локализация, SEO/Open Graph, sitemap и server-side формы", slide_no)

    add_bullet_card(
        slide,
        0.72,
        1.38,
        5.05,
        4.95,
        "Реализованные возможности",
        [
            "Locale-prefixed публичные маршруты `/ru/` и `/en/` для storefront-core.",
            "CMS-managed SEO/Open Graph для home-page, page, article, project и vacancy.",
            "Автоматическая генерация sitemap в build pipeline Astro.",
            "Формы лидов и откликов проходят через Astro server routes, а не пишут напрямую в Strapi из браузера.",
        ],
    )
    add_image(slide, assets["home"], 6.0, 1.45, 6.0)
    add_caption(slide, 6.08, 5.95, 5.8, "Публичная витрина отдает уже собранные маршруты и использует единый metadata pipeline.")

    add_label_chip(slide, 6.06, 6.38, "/ru / /en", fill=RGBColor(226, 235, 246))
    add_label_chip(slide, 7.58, 6.38, "meta + og:image")
    add_label_chip(slide, 9.58, 6.38, "sitemap build")
    add_label_chip(slide, 11.15, 6.38, "server-side forms")


def add_deployment_slide(prs: Presentation, assets: dict[str, Path], slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_band(slide, "Публикация, deployment и безопасность", "Как система замыкает сценарий publish -> rebuild и разграничивает доступ", slide_no)

    add_image(slide, assets["content_lifecycle"], 0.76, 1.36, 6.0)
    add_caption(slide, 0.88, 5.84, 5.7, "Публикация и снятие с публикации в Strapi инициируют managed webhook для rebuild-процесса.")

    add_bullet_card(
        slide,
        7.08,
        1.38,
        5.52,
        5.2,
        "Эксплуатационный контур",
        [
            "Versioned webhook в Strapi подписан на entry.publish и entry.unpublish.",
            "Frontend deployment ориентирован на Vercel deploy hook.",
            "CMS оформлена как Docker bundle с отдельным PostgreSQL runtime.",
            "Роли: administrator, marketer/content-manager, editor, HR.",
            "Draft-доступ выдается только через preview-secret; public API ограничен published-данными.",
        ],
    )


def add_results_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_band(slide, "Результаты и проверка", "Что реально подтверждено в проекте", slide_no)

    rows = [
        ("Frontend build и prerender", "Собирается 35 публичных HTML-маршрутов без runtime errors."),
        ("Preview contour", "Черновики доступны только через /api/preview и x-preview-secret."),
        ("SEO и sitemap", "Проверены title, canonical, og-поля и генерация sitemap-index.xml."),
        ("Форменные сценарии", "Lead-submission и vacancy-application валидируются на server-side и пишутся через технический токен."),
        ("Rebuild contour", "Managed webhook зарегистрирован в Strapi и привязан к событиям publish/unpublish."),
        ("Security baseline", "Public content API read-only; публикация отделена от редактирования ролями."),
    ]
    add_table_slide(slide, 0.75, 1.55, 12.0, 4.9, rows)

    add_text_block(
        slide,
        0.92,
        6.6,
        11.6,
        0.28,
        [{"text": "Результаты опираются на локальную сборку, smoke-checks, runtime-проверки и прямую верификацию CMS/SQLite baseline, а не только на чтение кода.", "size": 11, "color": MUTED, "align": PP_ALIGN.CENTER, "space_after": 0}],
    )


def add_significance_slide(prs: Presentation, slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_band(slide, "Практическая значимость и ограничения", None, slide_no)

    add_bullet_card(
        slide,
        0.74,
        1.45,
        5.85,
        5.15,
        "Практический эффект",
        [
            "Редактор может создавать и публиковать страницы, статьи, кейсы и вакансии без изменения frontend-кода.",
            "Контент, SEO и локализация переводятся в единый CMS-first контур.",
            "Публикация становится воспроизводимой благодаря webhook -> rebuild и versioned deployment bundle.",
        ],
        title_color=ACCENT_2,
    )
    add_bullet_card(
        slide,
        6.78,
        1.45,
        5.85,
        5.15,
        "Текущие ограничения",
        [
            "Полный browser-level WCAG audit и Lighthouse baseline не завершены.",
            "Внешний Vercel deploy не воспроизводился end-to-end внутри репозитория.",
            "Дальнейшее развитие: новые блоки Dynamic Zone, расширение workflow и аналитики контента.",
        ],
    )


def add_final_slide(prs: Presentation, assets: dict[str, Path], slide_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    hero.fill.solid()
    hero.fill.fore_color.rgb = ACCENT
    hero.line.fill.background()
    add_logo(slide, assets["logo"], 11.62, 0.15, 1.15)

    add_text_block(
        slide,
        0.86,
        0.28,
        9.8,
        0.38,
        [{"text": "Заключение", "size": 28, "bold": True, "font_name": "Aptos Display", "color": RGBColor(255, 255, 255), "space_after": 0}],
    )
    add_text_block(
        slide,
        0.98,
        1.65,
        11.2,
        2.1,
        [
            {"text": "Разработана CMS-first система управления контентом для маркетингового сайта.", "size": 22, "bold": True, "font_name": "Aptos Display", "space_after": 12},
            {"text": "Strapi используется как центр контентного управления, Astro — как публичная витрина с предсказуемой публикацией и build-based delivery.", "size": 15, "color": MUTED, "space_after": 8},
            {"text": "Работа закрывает не только архитектурное описание, но и практический инженерный контур: модель данных, preview, SEO, роли, rebuild и deployment baseline.", "size": 15, "color": MUTED, "space_after": 0},
        ],
    )

    add_label_chip(slide, 1.0, 4.25, "Strapi + Astro", fill=RGBColor(226, 235, 246))
    add_label_chip(slide, 2.98, 4.25, "CMS-first")
    add_label_chip(slide, 4.56, 4.25, "preview")
    add_label_chip(slide, 5.78, 4.25, "SEO / sitemap")
    add_label_chip(slide, 7.65, 4.25, "webhook -> rebuild")

    thanks = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(5.18), Inches(11.45), Inches(1.45))
    thanks.fill.solid()
    thanks.fill.fore_color.rgb = CARD
    thanks.line.color.rgb = LINE

    add_text_block(
        slide,
        1.08,
        5.43,
        11.1,
        0.8,
        [
            {"text": "Спасибо за внимание", "size": 25, "bold": True, "font_name": "Aptos Display", "align": PP_ALIGN.CENTER, "space_after": 8},
            {"text": f"Нахатакян Артур Романович · {slide_no} слайдов в основной части защиты", "size": 12, "color": MUTED, "align": PP_ALIGN.CENTER, "space_after": 0},
        ],
    )


def main() -> None:
    ensure_dirs()
    assets = build_assets()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, assets)
    add_relevance_slide(prs, assets, 2)
    add_goal_slide(prs, 3)
    add_stack_slide(prs, 4)
    add_architecture_slide(prs, assets, 5)
    add_data_model_slide(prs, assets, 6)
    add_editor_slide(prs, assets, 7)
    add_frontend_slide(prs, assets, 8)
    add_deployment_slide(prs, assets, 9)
    add_results_slide(prs, 10)
    add_significance_slide(prs, 11)
    add_final_slide(prs, assets, 12)

    for slide in prs.slides:
        remove_empty_placeholders(slide)

    prs.save(OUTPUT_PPTX)
    print(f"Saved {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
